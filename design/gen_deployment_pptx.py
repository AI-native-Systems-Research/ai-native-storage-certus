#!/usr/bin/env python3
"""Generate a PowerPoint slide showing the certus-server component deployment diagram."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# --- Colors ---
C_TITLE = RGBColor(0x33, 0x33, 0x33)
C_CLIENT_BG = RGBColor(0xE3, 0xF2, 0xFD)
C_CLIENT_BD = RGBColor(0x15, 0x65, 0xC0)
C_GRPC_BG = RGBColor(0xE8, 0xF5, 0xE9)
C_GRPC_BD = RGBColor(0x2E, 0x7D, 0x32)
C_DISP_BG = RGBColor(0xF3, 0xE5, 0xF5)
C_DISP_BD = RGBColor(0x6A, 0x1B, 0x9A)
C_DMAP_BG = RGBColor(0xE0, 0xF7, 0xFA)
C_DMAP_BD = RGBColor(0x00, 0x69, 0x5C)
C_GPU_BG = RGBColor(0xFF, 0xF9, 0xC4)
C_GPU_BD = RGBColor(0xF9, 0xA8, 0x25)
C_SPDK_BG = RGBColor(0xFB, 0xE9, 0xE7)
C_SPDK_BD = RGBColor(0xBF, 0x36, 0x0C)
C_EM_BG = RGBColor(0xFF, 0xF3, 0xE0)
C_EM_BD = RGBColor(0xE6, 0x51, 0x00)
C_LOG_BG = RGBColor(0xF5, 0xF5, 0xF5)
C_LOG_BD = RGBColor(0x61, 0x61, 0x61)
C_HW_BG = RGBColor(0xFF, 0xCC, 0xBC)
C_HW_BD = RGBColor(0xBF, 0x36, 0x0C)
C_ARROW = RGBColor(0x42, 0x42, 0x42)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_box(slide, left, top, width, height, text, fill_color, border_color, font_size=9, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold if i == 0 else False
        run.font.color.rgb = border_color

    return shape


def add_arrow(slide, x1, y1, x2, y2, label="", color=C_ARROW):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    # Add end arrow
    connector.end_x = x2
    connector.end_y = y2
    if label:
        # Place label near midpoint
        mid_x = (x1 + x2) // 2 - Inches(0.4)
        mid_y = (y1 + y2) // 2 - Inches(0.12)
        txBox = slide.shapes.add_textbox(mid_x, mid_y, Inches(1.2), Inches(0.25))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(7)
        run.font.color.rgb = color


# --- Title ---
title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(8), Inches(0.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Certus Server — Component Deployment"
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = C_TITLE

# Layout constants
COL1_X = Inches(0.3)   # Left column
COL2_X = Inches(4.8)   # Center column
COL3_X = Inches(9.3)   # Right column

ROW1_Y = Inches(0.7)   # Client
ROW2_Y = Inches(1.7)   # gRPC
ROW3_Y = Inches(2.7)   # Dispatcher
ROW4_Y = Inches(4.4)   # DispatchMap / GPU / Logger
ROW5_Y = Inches(5.6)   # Metadata EM + SPDK
ROW6_Y = Inches(6.6)   # Hardware

BOX_W = Inches(3.8)
BOX_H = Inches(0.8)
BOX_SM_W = Inches(2.8)
BOX_SM_H = Inches(0.7)

# === CLIENT ===
add_box(slide, COL2_X, ROW1_Y, BOX_W, BOX_H,
        "Python Test Client\nPyTorch GPU alloc • cudaIpcGetMemHandle • Batch gRPC",
        C_CLIENT_BG, C_CLIENT_BD, font_size=9, bold=True)

# === gRPC SERVICE ===
add_box(slide, COL2_X, ROW2_Y, BOX_W, BOX_H,
        "gRPC Service Layer (tonic)\nPopulate | Lookup | Check | Remove\ncudaIpcOpen/Close • batch→singular",
        C_GRPC_BG, C_GRPC_BD, font_size=8, bold=True)

# === DISPATCHER (large box) ===
disp_h = Inches(1.5)
add_box(slide, Inches(1.5), ROW3_Y, Inches(10.3), disp_h,
        "",  # we'll add sub-boxes inside
        C_DISP_BG, C_DISP_BD)

# Dispatcher label
disp_label = slide.shapes.add_textbox(Inches(1.7), ROW3_Y + Inches(0.05), Inches(5), Inches(0.3))
tf = disp_label.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "DispatcherComponentV0  «IDispatcher»"
run.font.size = Pt(10)
run.font.bold = True
run.font.color.rgb = C_DISP_BD

# Dispatcher sub-labels
sub_y = ROW3_Y + Inches(0.35)
add_box(slide, Inches(1.8), sub_y, Inches(2.5), Inches(0.55),
        "Dispatcher Core\npopulate • lookup • check • remove",
        C_DISP_BG, C_DISP_BD, font_size=8)

add_box(slide, Inches(4.6), sub_y, Inches(2.2), Inches(0.55),
        "BackgroundWriter\nasync staging→NVMe",
        RGBColor(0xED, 0xE7, 0xF6), RGBColor(0x45, 0x27, 0xA0), font_size=8)

# DataDrive inner box
dd_x = Inches(7.1)
dd_w = Inches(4.5)
dd_h = Inches(1.0)
add_box(slide, dd_x, sub_y - Inches(0.05), dd_w, dd_h,
        "",
        RGBColor(0xF8, 0xF0, 0xFC), RGBColor(0x45, 0x27, 0xA0))

dd_label = slide.shapes.add_textbox(dd_x + Inches(0.1), sub_y - Inches(0.02), Inches(3), Inches(0.22))
tf = dd_label.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "DataDrive[0..N] (per --data-pci)"
run.font.size = Pt(7)
run.font.italic = True
run.font.color.rgb = RGBColor(0x45, 0x27, 0xA0)

add_box(slide, dd_x + Inches(0.15), sub_y + Inches(0.25), Inches(2.0), Inches(0.55),
        "BlockDeviceSpdkNvmeV2\n«IBlockDevice»",
        C_SPDK_BG, C_SPDK_BD, font_size=7, bold=True)

add_box(slide, dd_x + Inches(2.3), sub_y + Inches(0.25), Inches(2.0), Inches(0.55),
        "ExtentManagerV2\n«IExtentManager»",
        C_EM_BG, C_EM_BD, font_size=7, bold=True)

# Receptacles label
recep = slide.shapes.add_textbox(Inches(1.8), sub_y + Inches(0.6), Inches(3), Inches(0.4))
tf = recep.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "receptacles: dispatch_map, gpu_services, spdk_env, logger"
run.font.size = Pt(7)
run.font.italic = True
run.font.color.rgb = C_DISP_BD

# === ROW 4: DispatchMap, GpuServices, Logger ===
add_box(slide, COL1_X, ROW4_Y, Inches(3.5), Inches(1.0),
        "DispatchMapComponentV0\n«IDispatchMap»\nEntry Table: key → {Staging | BlockDevice}\nreceptacles: extent_manager, logger",
        C_DMAP_BG, C_DMAP_BD, font_size=8, bold=True)

add_box(slide, Inches(4.2), ROW4_Y, Inches(3.2), Inches(1.0),
        "GpuServicesComponentV0\n«IGpuServices»\ncudaMemcpy (H2D / D2H)\nreceptacles: logger",
        C_GPU_BG, C_GPU_BD, font_size=8, bold=True)

add_box(slide, Inches(7.8), ROW4_Y, Inches(2.2), Inches(1.0),
        "LoggerComponentV1\n«ILogger»\nConsole + file logging",
        C_LOG_BG, C_LOG_BD, font_size=8, bold=True)

add_box(slide, Inches(10.4), ROW4_Y, Inches(2.6), Inches(1.0),
        "SPDKEnvComponent\n«ISPDKEnv»\nDPDK/EAL init\nVFIO device discovery",
        C_SPDK_BG, C_SPDK_BD, font_size=8, bold=True)

# === ROW 5: Metadata path ===
add_box(slide, COL1_X, ROW5_Y, Inches(3.0), Inches(0.8),
        "Metadata ExtentManagerV2\n«IExtentManager»\nreceptacles: metadata_device, logger",
        C_EM_BG, C_EM_BD, font_size=8, bold=True)

add_box(slide, Inches(3.7), ROW5_Y, Inches(3.2), Inches(0.8),
        "Metadata BlockDeviceSpdkNvmeV2\n«IBlockDevice, IBlockDeviceAdmin»\nreceptacles: spdk_env, logger",
        C_SPDK_BG, C_SPDK_BD, font_size=8, bold=True)

# === ROW 6: Hardware ===
add_box(slide, COL1_X, ROW6_Y, Inches(2.2), Inches(0.6),
        "NVMe (metadata)\n--metadata-pci",
        C_HW_BG, C_HW_BD, font_size=8, bold=True)

add_box(slide, Inches(3.0), ROW6_Y, Inches(2.2), Inches(0.6),
        "NVMe (data) [0..N]\n--data-pci",
        C_HW_BG, C_HW_BD, font_size=8, bold=True)

add_box(slide, Inches(5.8), ROW6_Y, Inches(2.2), Inches(0.6),
        "GPU (server context)\nCUDA device memory",
        C_GPU_BG, C_GPU_BD, font_size=8, bold=True)

add_box(slide, Inches(8.6), ROW6_Y, Inches(2.0), Inches(0.6),
        "GPU (client context)\nCUDA IPC shared",
        C_CLIENT_BG, C_CLIENT_BD, font_size=8, bold=True)

# === ARROWS (connections) ===
# Client → gRPC
add_arrow(slide,
          COL2_X + BOX_W // 2, ROW1_Y + BOX_H,
          COL2_X + BOX_W // 2, ROW2_Y,
          "gRPC", C_CLIENT_BD)

# gRPC → Dispatcher
add_arrow(slide,
          COL2_X + BOX_W // 2, ROW2_Y + BOX_H,
          COL2_X + BOX_W // 2, ROW3_Y,
          "", C_GRPC_BD)

# Dispatcher → DispatchMap
add_arrow(slide,
          Inches(2.5), ROW3_Y + disp_h,
          Inches(2.0), ROW4_Y,
          "bind", C_DISP_BD)

# Dispatcher → GPU
add_arrow(slide,
          Inches(5.5), ROW3_Y + disp_h,
          Inches(5.8), ROW4_Y,
          "bind", C_DISP_BD)

# Dispatcher → SPDK
add_arrow(slide,
          Inches(10.0), ROW3_Y + disp_h,
          Inches(11.5), ROW4_Y,
          "bind", C_DISP_BD)

# DispatchMap → Metadata EM
add_arrow(slide,
          Inches(2.0), ROW4_Y + Inches(1.0),
          Inches(1.8), ROW5_Y,
          "bind", C_DMAP_BD)

# Metadata EM → Metadata BD
add_arrow(slide,
          Inches(3.3), ROW5_Y + Inches(0.4),
          Inches(3.7), ROW5_Y + Inches(0.4),
          "", C_EM_BD)

# Metadata BD → NVMe meta
add_arrow(slide,
          Inches(4.5), ROW5_Y + Inches(0.8),
          Inches(1.3), ROW6_Y,
          "NVMe I/O", C_SPDK_BD)

# GPU → GPU hardware
add_arrow(slide,
          Inches(5.8), ROW4_Y + Inches(1.0),
          Inches(6.9), ROW6_Y,
          "cudaMemcpy", C_GPU_BD)

# SPDK → NVMe
add_arrow(slide,
          Inches(11.5), ROW4_Y + Inches(1.0),
          Inches(4.0), ROW6_Y,
          "VFIO", C_SPDK_BD)

# --- Data flow legend ---
legend_y = Inches(0.15)
legend_x = Inches(9.0)
legend_box = slide.shapes.add_textbox(legend_x, legend_y, Inches(4.0), Inches(0.5))
tf = legend_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Data flow: Client GPU →(IPC)→ Server GPU →(DMA)→ Staging →(SPDK)→ NVMe"
run.font.size = Pt(7)
run.font.italic = True
run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# Save
output_path = "/home/dwaddington/certus/design/certus-server-deployment.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
