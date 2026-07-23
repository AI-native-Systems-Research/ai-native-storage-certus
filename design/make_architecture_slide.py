#!/usr/bin/env python3
"""Build a single 16:9 PowerPoint slide of the Certus high-level architecture.

Places the rendered architecture diagram (certus-architecture-overview-slide.png,
produced from the matching .puml via `plantuml -tpng`) on the left, with the
seven key features as a crisp numbered list on the right.

Usage:
    plantuml -tpng certus-architecture-overview-slide.puml
    python3 make_architecture_slide.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
DIAGRAM = os.path.join(HERE, "certus-architecture-overview-slide.png")
OUT = os.path.join(HERE, "certus-architecture-overview.pptx")

# Palette (light, matches the PlantUML diagram)
INK = RGBColor(0x1F, 0x29, 0x37)      # near-black text
BLUE = RGBColor(0x1D, 0x4E, 0xD8)     # title / accent
SUBTLE = RGBColor(0x4B, 0x55, 0x63)   # descriptions
RULE = RGBColor(0x93, 0xC5, 0xFD)     # accent rule

FEATURES = [
    ("Lazy write-through",
     "Populate returns once the block is in the DRAM tier; SSD persistence "
     "happens asynchronously via the background writer."),
    ("Memory-tier-full handling",
     "When DRAM is saturated, LRU blocks with completed write-through demote "
     "to SSD-only; cold data is then served straight from SSD."),
    ("Pluggable block-device backends",
     "Any IBlockDevice impl — SPDK userspace NVMe, filesystem-backed, or "
     "kernel block device — paired with an ExtentManager."),
    ("Optional GPUDirect P2P",
     "The dispatcher-p2p variant streams SSD ↔ GPU through a BAR1-mapped "
     "ring, bypassing the DRAM staging copy (profile full-p2p)."),
    ("Cold-hit promotion",
     "A lookup that misses DRAM but hits SSD is pipelined into DRAM and "
     "re-registered as a warm memory-tier entry."),
    ("Pluggable eviction policy",
     "MemoryTier and DispatchMap delegate victim selection to an "
     "IEvictionPolicy plug-in (LRU today)."),
    ("Remote peers via RDMA",
     "In full-remote, local misses query peers (Zyre discovery); the holding "
     "peer RDMA-writes the value directly into local memory — one-sided, "
     "no CPU on the requester's critical path."),
]


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # White background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title
    title = slide.shapes.add_textbox(Inches(0.35), Inches(0.28), Inches(12.6), Inches(0.8))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Certus — High-Level Architecture & Key Features"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Accent rule under the title
    rule = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.38), Inches(1.06), Inches(12.55), Pt(2.5),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE
    rule.line.fill.background()
    rule.shadow.inherit = False

    # Diagram (left). Width fixed; height auto-scales to preserve aspect.
    slide.shapes.add_picture(DIAGRAM, Inches(0.25), Inches(1.35), width=Inches(7.75))

    # Feature list (right column)
    box = slide.shapes.add_textbox(Inches(8.2), Inches(1.3), Inches(4.95), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    head = tf.paragraphs[0]
    head.text = "Key features"
    head.font.size = Pt(17)
    head.font.bold = True
    head.font.color.rgb = INK
    head.space_after = Pt(6)

    for i, (name, desc) in enumerate(FEATURES, start=1):
        para = tf.add_paragraph()
        para.space_after = Pt(6)
        lead = para.add_run()
        lead.text = f"({i}) {name} — "
        lead.font.size = Pt(11.5)
        lead.font.bold = True
        lead.font.color.rgb = BLUE
        body = para.add_run()
        body.text = desc
        body.font.size = Pt(11)
        body.font.color.rgb = SUBTLE

    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
